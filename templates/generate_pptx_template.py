# -*- coding: utf-8 -*-
"""
会议总结 PPTX 生成器模板
基于 python-pptx，绝对坐标定位，关键概念栏固定在每页相同位置。

使用方法：
  1. 填写下方「配置区」的变量
  2. 填入 FOCUS_SLIDES 和 OTHER_SESSIONS 数据
  3. python generate_<会议简称>.py
"""

import os, glob, re
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ══════════════════════════════════════════════════════════════════════════════
# 配置区 — 每次新会议只需修改这里和下方的数据区
# ══════════════════════════════════════════════════════════════════════════════

FOCUS_DIR = r'...\Library\...\Focus'   # Focus 照片目录
OTHER_DIR = r'...\Library\...\Other'   # Other 照片目录（若无 Other 可设为 None）
OUTPUT    = r'...\会议名_年份_会议总结.pptx'

CONF_SUBTITLE   = 'CMAC 生物医药新技术创新亦庄论坛'   # 会议副标题（封面用）
CONF_DATE_VENUE = '2026年6月25日  ·  北京经济技术开发区亦庄'
CONF_EN_VENUE   = 'Beijing Economic and Technological Development Zone'
COVER_BAR_LINE1 = '会议总结  ·  内部分享'          # 封面关键概念栏第一行
COVER_BAR_LINE2 = 'AI助力RWE  ·  医药创新  ·  临床统计  ·  大模型'  # 第二行（大字）

FOCUS_LABEL    = 'FOCUS 重点报告'
FOCUS_SPEAKER  = '讲者姓名'
FOCUS_AFFIL    = '单位'
FOCUS_TITLE_ZH = '演讲题目（中文）'
FOCUS_TITLE_EN = 'Lecture Title (English)'
FOCUS_NOTE     = ''   # 自动填写，不用手动设置

OTHER_TITLE_ZH  = '生物医药 · AI · 临床统计 · 大模型'
OTHER_TITLE_EN  = 'Pharmaceutical Innovation · AI · Clinical Statistics · LLM'
OTHER_SECTION_SPEAKER = '会议报告综述'
OTHER_SECTION_AFFIL   = '以下讲者关键信息摘要'

BAR_FOOTER = 'AI助力RWE  ·  讲者姓名  ·  CMAC  ·  2026-06-25'  # 每页关键概念栏页脚

# ══════════════════════════════════════════════════════════════════════════════
# 数据区 — 由视觉 AI 自动提取后填入
# ══════════════════════════════════════════════════════════════════════════════

# 格式: (照片编号, 幻灯片标题, [关键词列表], [要点列表])
# 要点中，以两个空格开头的行会以较小字体显示（子要点）
# 空行请用 '' 表示
FOCUS_SLIDES = [
    # (1, '幻灯片标题', ['关键词1', '关键词2', '关键词3'],
    #  ['要点一，一到两句话',
    #   '要点二',
    #   '  → 子要点（以两个空格开头）']),
]

# 格式: 每个 Other 讲者一个 dict
OTHER_SESSIONS = [
    # {
    #     'speaker': '讲者姓名',
    #     'affiliation': '单位职称',
    #     'title': '演讲题目中文',
    #     'en_title': 'Title in English',
    #     'photo_nums': [代表性照片编号1, 照片编号2, 照片编号3],
    #     'points': [
    #         '核心观点一',
    #         '  → 子观点（缩进两空格）',
    #         '',
    #         '核心观点二',
    #     ],
    # },
]

# ══════════════════════════════════════════════════════════════════════════════
# 以下代码通常不需要修改
# ══════════════════════════════════════════════════════════════════════════════

DARK_BLUE  = RGBColor(0x0D, 0x3B, 0x6E)
MID_BLUE   = RGBColor(0x1A, 0x6A, 0xB8)
LIGHT_BLUE = RGBColor(0x8C, 0xB8, 0xE0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
MUTED      = RGBColor(0x66, 0x77, 0x88)
GREY_LINE  = RGBColor(0xCC, 0xD4, 0xDE)

SW, SH = 29.70, 21.00
ML = 1.20
PHOTO_L  = ML
PHOTO_T  = 1.41
PHOTO_W  = 20.80
PHOTO_H  = 15.60
GAP      = 0.15
RIGHT_L  = PHOTO_L + PHOTO_W + GAP
RIGHT_W  = SW - RIGHT_L - ML
BAR_L    = ML
BAR_T    = 17.61   # FIXED — do not change
BAR_W    = SW - 2 * ML
BAR_H    = 2.96
HDR_T    = 1.49
HDR_H    = 2.30
DIV_T    = HDR_T + HDR_H + 0.15
BODY_T   = DIV_T + 0.20
BODY_BOT = BAR_T - 0.20
BODY_H   = BODY_BOT - BODY_T


def find_photo(num, folder=None):
    dirs = [d for d in [folder, FOCUS_DIR, OTHER_DIR if OTHER_DIR else None] if d]
    for d in dirs:
        if not os.path.isdir(d):
            continue
        hits = glob.glob(os.path.join(d, f'*_{num}_*.jpg'))
        if not hits:
            # fallback: any file whose last numeric group == num
            for f in os.listdir(d):
                nums = re.findall(r'\d+', f)
                if nums and int(nums[-1]) == num:
                    hits = [os.path.join(d, f)]
                    break
        if hits:
            return hits[0]
    return None


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, l, t, w, h, color, line_color=None):
    shape = slide.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    return shape


def add_textbox(slide, l, t, w, h, word_wrap=True):
    txb = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    txb.text_frame.word_wrap = word_wrap
    return txb


def para(tf, text, font_name='Microsoft YaHei', size=12, bold=False,
         color=None, align=PP_ALIGN.LEFT, space_before=0, space_after=0,
         line_spacing=None):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return p


def add_photo(slide, path, l, t, max_w, max_h):
    from PIL import Image as PILImage
    with PILImage.open(path) as im:
        iw, ih = im.size
    aspect = iw / ih
    if (max_w / max_h) > aspect:
        fh = max_h
        fw = max_h * aspect
    else:
        fw = max_w
        fh = max_w / aspect
    off_l = l + (max_w - fw) / 2
    off_t = t + (max_h - fh) / 2
    slide.shapes.add_picture(path, Cm(off_l), Cm(off_t), Cm(fw), Cm(fh))


def build_cover(prs):
    slide = blank_slide(prs)
    add_rect(slide, ML, 1.41, SW - 2*ML, SH - 1.41 - 0.43, DARK_BLUE)

    txb = add_textbox(slide, ML + 1, 5.0, SW - 2*ML - 2, 3.5)
    tf = txb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run()
    r.text = CONF_SUBTITLE
    r.font.name = 'Microsoft YaHei'
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE

    para(tf, CONF_DATE_VENUE, size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, space_before=12)
    para(tf, CONF_EN_VENUE,   size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, space_before=4)

    add_rect(slide, ML, BAR_T, SW - 2*ML, BAR_H, DARK_BLUE)
    bt = add_textbox(slide, ML + 0.6, BAR_T + 0.12, SW - 2*ML - 1, BAR_H - 0.24)
    tf2 = bt.text_frame
    tf2.word_wrap = True
    p1 = tf2.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = COVER_BAR_LINE1
    r1.font.name = 'Microsoft YaHei'
    r1.font.size = Pt(14)
    r1.font.color.rgb = LIGHT_BLUE
    para(tf2, COVER_BAR_LINE2, size=22, bold=True, color=WHITE, space_before=4)


def build_section_divider(prs, section_label, speaker_name, affiliation,
                           title_zh, title_en, subtitle=''):
    slide = blank_slide(prs)
    add_rect(slide, ML, 1.41, SW - 2*ML, SH - 1.41 - 0.43, DARK_BLUE)

    txb = add_textbox(slide, ML + 1, 4.5, SW - 2*ML - 2, 10)
    tf = txb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run()
    r.text = section_label
    r.font.name = 'Microsoft YaHei'
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE

    para(tf, speaker_name,  size=22, color=WHITE,      align=PP_ALIGN.CENTER, space_before=10)
    para(tf, affiliation,   size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, space_before=4)
    para(tf, title_zh,      size=16, color=WHITE,      align=PP_ALIGN.CENTER, space_before=16)
    if title_en:
        para(tf, title_en,  size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, space_before=4)
    if subtitle:
        para(tf, subtitle,  size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, space_before=8)


def build_focus_slide(prs, idx, num, title, tags, bullets):
    slide = blank_slide(prs)

    photo_path = find_photo(num, FOCUS_DIR)
    if photo_path:
        slide.shapes.add_picture(photo_path,
                                  Cm(PHOTO_L), Cm(PHOTO_T),
                                  Cm(PHOTO_W), Cm(PHOTO_H))
    else:
        add_rect(slide, PHOTO_L, PHOTO_T, PHOTO_W, PHOTO_H, RGBColor(0xDD, 0xDD, 0xDD))

    hdr = add_textbox(slide, RIGHT_L, HDR_T, RIGHT_W, HDR_H)
    tf = hdr.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r = p0.add_run()
    r.text = title
    r.font.name = 'Microsoft YaHei'
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = DARK_BLUE
    para(tf, f'幻灯片 {idx + 1} / {len(FOCUS_SLIDES)}  ·  照片 #{num}',
         size=10, color=MUTED, space_before=4)

    add_rect(slide, RIGHT_L + 0.3, DIV_T, RIGHT_W - 0.5, 0.03, GREY_LINE)

    btxb = add_textbox(slide, RIGHT_L, BODY_T, RIGHT_W, BODY_H)
    tf2 = btxb.text_frame
    tf2.word_wrap = True

    first = True
    for line in bullets:
        if line == '':
            p = tf2.add_paragraph() if not first else tf2.paragraphs[0]
            if first:
                p._p.clear(); first = False
            p.space_after = Pt(4)
        elif line.startswith('  ') or line.startswith('\t'):
            p = tf2.add_paragraph() if not first else tf2.paragraphs[0]
            if first:
                p._p.clear(); first = False
            p.space_before = Pt(2)
            r = p.add_run()
            r.text = line.strip()
            r.font.name = 'Microsoft YaHei'
            r.font.size = Pt(10)
            r.font.color.rgb = MUTED
        else:
            p = tf2.add_paragraph() if not first else tf2.paragraphs[0]
            if first:
                p._p.clear(); first = False
            p.space_before = Pt(6)
            p.space_after = Pt(3)
            r = p.add_run()
            r.text = line
            r.font.name = 'Microsoft YaHei'
            r.font.size = Pt(11)
            r.font.color.rgb = TEXT_DARK

    # Keyword bar — always at fixed position BAR_T
    add_rect(slide, BAR_L, BAR_T, BAR_W, BAR_H, DARK_BLUE)
    bar_txb = add_textbox(slide, BAR_L + 0.5, BAR_T + 0.12, BAR_W - 1, BAR_H - 0.24)
    tf3 = bar_txb.text_frame
    tf3.word_wrap = False
    p_lbl = tf3.paragraphs[0]
    p_lbl.alignment = PP_ALIGN.LEFT
    r_lbl = p_lbl.add_run()
    r_lbl.text = '关键概念'
    r_lbl.font.name = 'Microsoft YaHei'
    r_lbl.font.size = Pt(10)
    r_lbl.font.color.rgb = LIGHT_BLUE
    para(tf3, '  ·  '.join(tags), size=18, bold=True, color=WHITE, space_before=2, line_spacing=22)
    para(tf3, BAR_FOOTER, size=9, color=LIGHT_BLUE, space_before=2)


def build_other_slide(prs, session):
    slide = blank_slide(prs)
    speaker    = session['speaker']
    affil      = session['affiliation']
    title      = session['title']
    en_title   = session.get('en_title', '')
    photo_nums = session.get('photo_nums', session.get('slide_nums', []))
    points     = session['points']

    add_rect(slide, ML, 0.04, SW - 2*ML, 4.15, DARK_BLUE)
    hdr_txb = add_textbox(slide, ML + 0.5, 0.18, SW - 2*ML - 1, 3.8)
    tf = hdr_txb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    r = p0.add_run()
    r.text = speaker
    r.font.name = 'Microsoft YaHei'
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = WHITE
    para(tf, affil,    size=13, color=LIGHT_BLUE, space_before=4)
    para(tf, title,    size=14, bold=True, color=WHITE, space_before=6)
    if en_title:
        para(tf, en_title, size=10, color=LIGHT_BLUE, space_before=3)

    btxb = add_textbox(slide, ML + 0.1, 4.40, 15.50, SH - 4.40 - 0.50)
    tf2 = btxb.text_frame
    tf2.word_wrap = True
    first = True
    for line in points:
        if line == '':
            p = tf2.add_paragraph() if not first else tf2.paragraphs[0]
            if first:
                p._p.clear(); first = False
            p.space_after = Pt(3)
        elif line.startswith('  ') or line.startswith('\t'):
            p = tf2.add_paragraph() if not first else tf2.paragraphs[0]
            if first:
                p._p.clear(); first = False
            r = p.add_run()
            r.text = line.strip()
            r.font.name = 'Microsoft YaHei'
            r.font.size = Pt(10)
            r.font.color.rgb = MUTED
            p.space_before = Pt(1)
        else:
            p = tf2.add_paragraph() if not first else tf2.paragraphs[0]
            if first:
                p._p.clear(); first = False
            r = p.add_run()
            r.text = line
            r.font.name = 'Microsoft YaHei'
            r.font.size = Pt(11)
            r.font.color.rgb = TEXT_DARK
            p.space_before = Pt(5)
            p.space_after = Pt(2)

    photo_l = 18.70
    photo_w = SW - photo_l - ML
    avail_h = SH - 4.40 - 0.50
    ph_count = min(len(photo_nums), 3)
    if ph_count > 0:
        ph_h = avail_h / ph_count - 0.15
        for i, pnum in enumerate(photo_nums[:ph_count]):
            ppath = find_photo(pnum, OTHER_DIR) or find_photo(pnum, FOCUS_DIR)
            if ppath:
                t = 4.40 + i * (ph_h + 0.15)
                add_photo(slide, ppath, photo_l, t, photo_w, ph_h)


def build():
    global FOCUS_NOTE
    FOCUS_NOTE = FOCUS_NOTE or f'以下 {len(FOCUS_SLIDES)} 张幻灯片逐页解读'

    prs = Presentation()
    prs.slide_width  = Cm(SW)
    prs.slide_height = Cm(SH)

    print('封面...')
    build_cover(prs)

    if FOCUS_SLIDES:
        print('Focus 分节页...')
        build_section_divider(prs, FOCUS_LABEL, FOCUS_SPEAKER, FOCUS_AFFIL,
                              FOCUS_TITLE_ZH, FOCUS_TITLE_EN, FOCUS_NOTE)
        print(f'Focus 幻灯片 ({len(FOCUS_SLIDES)} 张)...')
        for idx, (num, title, tags, bullets) in enumerate(FOCUS_SLIDES):
            if (idx + 1) % 10 == 0:
                print(f'  ... {idx+1}/{len(FOCUS_SLIDES)}')
            build_focus_slide(prs, idx, num, title, tags, bullets)

    if OTHER_SESSIONS:
        print('Other 分节页...')
        build_section_divider(prs, 'OTHER 其他报告', OTHER_SECTION_SPEAKER,
                              OTHER_SECTION_AFFIL, OTHER_TITLE_ZH, OTHER_TITLE_EN)
        print('Other 讲者幻灯片...')
        for sess in OTHER_SESSIONS:
            print(f"  ... {sess['speaker']}")
            build_other_slide(prs, sess)

    print(f'保存到 {OUTPUT} ...')
    prs.save(OUTPUT)
    print(f'完成！共 {len(prs.slides)} 张幻灯片')


if __name__ == '__main__':
    build()
